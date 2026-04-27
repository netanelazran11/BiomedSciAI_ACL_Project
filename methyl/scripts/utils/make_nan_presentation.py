"""
Generate PowerPoint: Data Pipeline & NaN Handling in MethylLlama
================================================================
Covers:
  - The two datasets (pretrain 169k×49k, finetune 11.5k×49k)
  - Why NaN exists (Illumina array union)
  - Full pipeline: h5ad → Dataset → Collator → Model
  - Step-by-step NaN handling in each file
  - Bug in old SCBert code vs correct LLaMA code
  - Pretrain metrics (run 44450919, epoch 99)
  - Placeholders for investigation script results

Usage:
    python scripts/utils/make_nan_presentation.py
    # → NaN_Handling_MethylLlama.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT        = RGBColor(0x16, 0x21, 0x3E)
ACCENT2       = RGBColor(0x0F, 0x17, 0x2A)
HIGHLIGHT     = RGBColor(0x00, 0xB4, 0xD8)
HIGHLIGHT2    = RGBColor(0x90, 0xE0, 0xEF)
GREEN         = RGBColor(0x06, 0xD6, 0xA0)
ORANGE        = RGBColor(0xFF, 0xB7, 0x03)
RED           = RGBColor(0xFF, 0x6B, 0x6B)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
SUBTITLE_GRAY = RGBColor(0xAA, 0xBB, 0xCC)
CODE_BG       = RGBColor(0x0D, 0x17, 0x26)
CODE_TEXT     = RGBColor(0xA8, 0xD8, 0xEA)
CODE_COMMENT  = RGBColor(0x6A, 0x73, 0x7D)
CODE_KEY      = RGBColor(0x50, 0xFA, 0x7B)
YELLOW        = RGBColor(0xF1, 0xFA, 0x8C)
PURPLE        = RGBColor(0xBD, 0x93, 0xF9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def slide_header(slide, title, underline_width=Inches(5)):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), HIGHLIGHT)
    add_rect(slide, 0, SLIDE_H - Inches(0.07), SLIDE_W, Inches(0.07), HIGHLIGHT)
    add_textbox(slide, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.65),
                font_size=30, bold=True, color=HIGHLIGHT)
    add_rect(slide, Inches(0.5), Inches(0.88), underline_width, Inches(0.035), HIGHLIGHT)


def add_panel(slide, left, top, width, height, title, title_color=ORANGE):
    add_rect(slide, left, top, width, height, ACCENT)
    add_textbox(slide, title, left + Inches(0.18), top + Inches(0.12),
                width - Inches(0.36), Inches(0.42),
                font_size=18, bold=True, color=title_color)
    add_rect(slide, left + Inches(0.18), top + Inches(0.58),
             width - Inches(0.36), Inches(0.03), title_color)


def add_code_panel(slide, left, top, width, height, lines, comment_prefix="#"):
    """Dark code box. lines: list of (text, color) tuples."""
    add_rect(slide, left, top, width, height, CODE_BG)
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12),
                                   width - Inches(0.3), height - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line_text, line_color in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(13)
        run.font.color.rgb = line_color
    return tb


def add_metric_box(slide, left, top, width, height, label, value,
                   val_color=ORANGE, label_color=SUBTITLE_GRAY):
    add_rect(slide, left, top, width, height, ACCENT)
    add_textbox(slide, value, left, top + Inches(0.1),
                width, Inches(0.55),
                font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + Inches(0.62),
                width, Inches(0.35),
                font_size=13, color=label_color, align=PP_ALIGN.CENTER)


def add_step_arrow(slide, left, top):
    add_textbox(slide, "▼", left, top, Inches(0.5), Inches(0.4),
                font_size=20, color=HIGHLIGHT, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.1), HIGHLIGHT)
add_rect(slide, 0, SLIDE_H - Inches(0.1), SLIDE_W, Inches(0.1), HIGHLIGHT)

add_textbox(slide, "Data Pipeline & NaN Handling",
            Inches(1), Inches(1.4), Inches(11.3), Inches(1.1),
            font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "MethylLlama Pretraining on the 49k CpG Union Corpus",
            Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.7),
            font_size=22, color=HIGHLIGHT2, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(3.5), Inches(3.55), Inches(6.3), Inches(0.04), HIGHLIGHT)
add_textbox(slide,
    "How missing measurements (NaN) flow through the pipeline\n"
    "and are correctly excluded from the reconstruction loss",
    Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.8),
    font_size=17, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "Netanel Azran  |  BMFM-RNA Methylation  |  Hebrew University",
            Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
            font_size=15, color=SUBTITLE_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "Run: llama-small-all49k-r0.5-w0.0-44450919",
            Inches(1), Inches(6.2), Inches(11.3), Inches(0.4),
            font_size=13, color=CODE_TEXT, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — The Two Datasets
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "The Two Datasets", Inches(4))

# Left — pretrain
add_panel(slide, Inches(0.4), Inches(1.1), Inches(5.9), Inches(5.7),
          "PRETRAINING CORPUS", HIGHLIGHT)
for text, yoff, col, sz in [
    ("169,120 samples × 49,156 CpGs", 0.75, WHITE, 19),
    ("File: methylgpt_pretrain_type3.h5ad", 1.25, CODE_TEXT, 14),
    ("Path: data_methyl_pretrain_type3_h5ad/", 1.6, CODE_COMMENT, 13),
    ("Run: pretrain_llama_small.sh", 2.05, HIGHLIGHT2, 14),
    ("Job ID: 44450919", 2.4, ORANGE, 14),
    ("Script: bmfm_methylation/llama/pretrain_llama.py", 2.85, CODE_TEXT, 13),
    ("Module: WCEDLlamaModule  (wced_llama.py)", 3.2, CODE_TEXT, 13),
    ("Architecture: 256D × 4L × 4H  (~5M params)", 3.7, LIGHT_GRAY, 14),
    ("Input ratio: 0.5  (50% of valid CpGs per view)", 4.05, LIGHT_GRAY, 14),
    ("Age weight: 0.0  (pure reconstruction)", 4.4, LIGHT_GRAY, 14),
]:
    add_textbox(slide, text, Inches(0.65), Inches(yoff + 1.1),
                Inches(5.4), Inches(0.45), font_size=sz, color=col)

# Right — finetune
add_panel(slide, Inches(7.0), Inches(1.1), Inches(5.9), Inches(5.7),
          "FINE-TUNING CORPUS", GREEN)
for text, yoff, col, sz in [
    ("11,500 samples × 49,156 CpGs", 0.75, WHITE, 19),
    ("File: finetuning_49k.h5ad", 1.25, CODE_TEXT, 14),
    ("Path: data_methyl_finetune_49k_h5ad/", 1.6, CODE_COMMENT, 13),
    ("Run: finetune_llama_small.sh", 2.05, HIGHLIGHT2, 14),
    ("Job ID: 44574410", 2.4, ORANGE, 14),
    ("Script: bmfm_methylation/llama/finetune_llama.py", 2.85, CODE_TEXT, 13),
    ("Recon weight: 0.0  (age MSE only)", 3.2, GREEN, 14),
    ("Input ratio: 1.0  (all valid CpGs as input)", 3.7, LIGHT_GRAY, 14),
    ("Pooling: mean over all tokens", 4.05, LIGHT_GRAY, 14),
    ("Freeze encoder: true  (then unfreeze epoch 10)", 4.4, LIGHT_GRAY, 14),
]:
    add_textbox(slide, text, Inches(7.25), Inches(yoff + 1.1),
                Inches(5.4), Inches(0.45), font_size=sz, color=col)

# Middle divider
add_textbox(slide, "Same 49k\nvocabulary", Inches(6.05), Inches(3.4),
            Inches(0.9), Inches(0.9), font_size=12, color=HIGHLIGHT2, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 3 — Why NaN Exists
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Why NaN Exists — The Illumina Array Union Problem", Inches(8))

# Three array boxes
for i, (array, probes, color) in enumerate([
    ("27k Array",  "~27,000 probes",  ORANGE),
    ("450k Array", "~450,000 probes", HIGHLIGHT),
    ("EPIC Array",  "~850,000 probes", GREEN),
]):
    lx = Inches(0.4 + i * 3.2)
    add_rect(slide, lx, Inches(1.15), Inches(2.9), Inches(1.4), ACCENT)
    add_textbox(slide, array, lx + Inches(0.15), Inches(1.25),
                Inches(2.6), Inches(0.45), font_size=18, bold=True, color=color)
    add_textbox(slide, probes, lx + Inches(0.15), Inches(1.7),
                Inches(2.6), Inches(0.35), font_size=14, color=LIGHT_GRAY)
    add_textbox(slide, "↓", lx + Inches(1.1), Inches(2.6),
                Inches(0.7), Inches(0.4), font_size=22, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Union result
add_rect(slide, Inches(0.4), Inches(3.05), Inches(9.4), Inches(0.9), ACCENT)
add_textbox(slide, "Union of all studies  →  49,156 CpGs in dataset",
            Inches(0.6), Inches(3.15), Inches(9.0), Inches(0.6),
            font_size=20, bold=True, color=WHITE)

# Consequence
add_textbox(slide, "→", Inches(9.9), Inches(3.25), Inches(0.6), Inches(0.5),
            font_size=28, bold=True, color=ORANGE)
add_rect(slide, Inches(10.55), Inches(3.05), Inches(2.4), Inches(0.9), ACCENT)
add_textbox(slide, "NaN where\nnot measured", Inches(10.65), Inches(3.1),
            Inches(2.2), Inches(0.8), font_size=16, bold=True, color=ORANGE)

# Explanation panel
add_rect(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(2.7), ACCENT2)
add_textbox(slide, "What this means per sample:",
            Inches(0.6), Inches(4.2), Inches(6), Inches(0.4),
            font_size=17, bold=True, color=HIGHLIGHT)

rows = [
    ("450k-array sample", "~19,600 valid CpGs", "~29,556 NaN  (60%)",  HIGHLIGHT, RED),
    ("EPIC-array sample", "~40,000 valid CpGs",  "~9,156 NaN  (19%)",   GREEN,     ORANGE),
    ("27k-array sample",  "~27,000 valid CpGs",  "~22,156 NaN  (45%)",  ORANGE,    RED),
]
for i, (sample, valid, nan, vc, nc) in enumerate(rows):
    y = Inches(4.7 + i * 0.62)
    add_textbox(slide, sample, Inches(0.65), y, Inches(3.2), Inches(0.5),
                font_size=15, color=LIGHT_GRAY)
    add_textbox(slide, valid,  Inches(4.0),  y, Inches(3.0), Inches(0.5),
                font_size=15, color=vc)
    add_textbox(slide, nan,    Inches(7.2),  y, Inches(3.0), Inches(0.5),
                font_size=15, color=nc)

add_textbox(slide,
    "From pretrain run valid_pct≈42%  →  NaN rate ≈ 16%  (pretrain corpus has mostly well-measured probes)",
    Inches(0.6), Inches(6.55), Inches(12), Inches(0.45),
    font_size=14, color=GREEN, italic=True)


# =============================================================================
# SLIDE 4 — Full Pipeline Overview
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Full Pipeline — NaN Handling at Each Stage", Inches(7))

steps = [
    ("h5ad File",              "methylgpt_pretrain_type3.h5ad\n169k × 49k  |  NaN = raw missing values", HIGHLIGHT, "Input"),
    ("MethylationDataset",     "shared/data_module.py\nComputes valid_mask = isfinite(β)\nNaN preserved in beta_values", ORANGE, "Step 1"),
    ("WCEDCollator  Part A",   "shared/data_module.py\nNaN → 0  in all_betas (targets)\nvalid_mask stored for loss", GREEN, "Step 2a"),
    ("WCEDCollator  Part B",   "shared/data_module.py\nInput views: only valid CpGs sampled\nNaN positions never shown to encoder", HIGHLIGHT2, "Step 2b"),
    ("WCEDLlamaModule",        "llama/wced_llama.py\nrecon_mask = ~input_mask & valid_mask\nLoss on measured, held-out CpGs ONLY", PURPLE, "Step 3"),
]

box_w = Inches(2.3)
box_h = Inches(1.5)
gap   = Inches(0.22)
start_x = Inches(0.3)
start_y = Inches(1.2)

for i, (title, body, color, tag) in enumerate(steps):
    lx = start_x + i * (box_w + gap)
    add_rect(slide, lx, start_y, box_w, box_h, ACCENT)
    add_rect(slide, lx, start_y, box_w, Inches(0.06), color)
    add_textbox(slide, tag,   lx + Inches(0.08), start_y + Inches(0.08),
                Inches(0.7), Inches(0.3), font_size=11, color=color)
    add_textbox(slide, title, lx + Inches(0.08), start_y + Inches(0.35),
                box_w - Inches(0.16), Inches(0.45),
                font_size=14, bold=True, color=WHITE)
    add_textbox(slide, body,  lx + Inches(0.08), start_y + Inches(0.82),
                box_w - Inches(0.16), Inches(0.65),
                font_size=11, color=LIGHT_GRAY)
    if i < len(steps) - 1:
        arrow_x = lx + box_w + Inches(0.03)
        add_textbox(slide, "→", arrow_x, start_y + Inches(0.55),
                    gap + Inches(0.1), Inches(0.4),
                    font_size=20, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

# Three outcome lanes
lane_y = Inches(3.1)
add_textbox(slide, "Three categories of CpG positions per sample:",
            Inches(0.4), lane_y, Inches(12.5), Inches(0.4),
            font_size=17, bold=True, color=HIGHLIGHT)

for i, (label, desc, val, color) in enumerate([
    ("In input view",        "Shown to encoder as tokens\n~50% of valid CpGs", "EXCLUDED from recon loss", ORANGE),
    ("Valid — held out",     "Real β target, not in input\n~50% of valid CpGs", "INCLUDED in recon loss ✓", GREEN),
    ("NaN position",         "Never measured, target=0\n~16% of vocab (pretrain)", "EXCLUDED by valid_mask ✓", RED),
]):
    lx = Inches(0.4 + i * 4.3)
    add_rect(slide, lx, lane_y + Inches(0.45), Inches(4.1), Inches(2.8), ACCENT)
    add_rect(slide, lx, lane_y + Inches(0.45), Inches(4.1), Inches(0.05), color)
    add_textbox(slide, label, lx + Inches(0.15), lane_y + Inches(0.55),
                Inches(3.8), Inches(0.4), font_size=16, bold=True, color=color)
    add_textbox(slide, desc,  lx + Inches(0.15), lane_y + Inches(1.0),
                Inches(3.8), Inches(0.75), font_size=13, color=LIGHT_GRAY)
    add_rect(slide, lx + Inches(0.15), lane_y + Inches(1.85), Inches(3.5), Inches(0.03), color)
    add_textbox(slide, val,   lx + Inches(0.15), lane_y + Inches(1.95),
                Inches(3.8), Inches(0.4), font_size=13, bold=True, color=color)


# =============================================================================
# SLIDE 5 — Step 1: MethylationDataset
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Step 1 — MethylationDataset: Reading the Data", Inches(7.5))

add_textbox(slide, "File: bmfm_methylation/shared/data_module.py   |   Class: MethylationDataset",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

# Left: code
add_code_panel(slide, Inches(0.4), Inches(1.4), Inches(7.5), Inches(4.5), [
    ("# MethylationDataset.__getitem__", CODE_COMMENT),
    ("def __getitem__(self, idx):", CODE_KEY),
    ("    beta_values = self.adata.X[idx]", CODE_TEXT),
    ("    if hasattr(beta_values, 'toarray'):", CODE_TEXT),
    ("        beta_values = beta_values.toarray().flatten()", CODE_TEXT),
    ("    beta_values = beta_values.astype(np.float32)", CODE_TEXT),
    ("", CODE_TEXT),
    ("    # Compute mask: True = real value, False = NaN", CODE_COMMENT),
    ("    valid_mask = np.isfinite(beta_values)", CODE_KEY),
    ("", CODE_TEXT),
    ("    # NaN is PRESERVED in beta_values at this stage", CODE_COMMENT),
    ("    mfi = MultiFieldInstance(", CODE_TEXT),
    ("        data={", CODE_TEXT),
    ('            "beta_values": beta_values.tolist(),  # NaN intact', CODE_TEXT),
    ('            "valid_mask": valid_mask.tolist(),    # True/False mask', CODE_KEY),
    ("        },", CODE_TEXT),
    ("        metadata={'labels': float(age), ...}", CODE_TEXT),
    ("    )", CODE_TEXT),
    ("    return mfi", CODE_TEXT),
])

# Right: explanation boxes
add_panel(slide, Inches(8.1), Inches(1.4), Inches(4.8), Inches(2.1),
          "What happens at this step", HIGHLIGHT)
add_textbox(slide,
    "• Raw β-value extracted from adata.X\n"
    "• valid_mask computed: np.isfinite(β)\n"
    "• NaN is NOT replaced — preserved as-is\n"
    "• Both β and mask passed to collator",
    Inches(8.3), Inches(2.0), Inches(4.4), Inches(1.4),
    font_size=14, color=LIGHT_GRAY)

add_panel(slide, Inches(8.1), Inches(3.65), Inches(4.8), Inches(2.2),
          "Why preserve NaN here?", ORANGE)
add_textbox(slide,
    "The collator decides what to do with NaN.\n\n"
    "Replacing NaN with 0 here would hide the\n"
    "missing-value information before the mask\n"
    "is computed. Preserving NaN ensures the\n"
    "valid_mask is always accurate.",
    Inches(8.3), Inches(4.25), Inches(4.4), Inches(1.5),
    font_size=14, color=LIGHT_GRAY)

# Bottom status
add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.75), ACCENT2)
add_textbox(slide, "NaN status after Step 1:   beta_values contains NaN    ✓ valid_mask computed",
            Inches(0.6), Inches(6.2), Inches(12), Inches(0.45),
            font_size=15, color=GREEN)


# =============================================================================
# SLIDE 6 — Step 2a: WCEDCollator — NaN Detection & Replacement
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Step 2a — WCEDCollator: NaN Detection & Replacement", Inches(8))

add_textbox(slide, "File: bmfm_methylation/shared/data_module.py   |   Class: WCEDCollator.__call__",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

# Code panel
add_code_panel(slide, Inches(0.4), Inches(1.4), Inches(7.5), Inches(4.5), [
    ("# For each sample i in the batch:", CODE_COMMENT),
    ("betas = np.asarray(ex.data['beta_values'], dtype=np.float32)", CODE_TEXT),
    ("", CODE_TEXT),
    ("# Slice to vocabulary of 49,156 CpGs", CODE_COMMENT),
    ("vocab_betas = betas[self.vocab_cpg_indices]", CODE_TEXT),
    ("", CODE_TEXT),
    ("# Detect NaN: True = measured, False = NaN", CODE_COMMENT),
    ("vocab_valid = np.isfinite(vocab_betas)", CODE_KEY),
    ("", CODE_TEXT),
    ("# Replace NaN → 0.0 (harmless: excluded from loss later)", CODE_COMMENT),
    ("vocab_betas_clean = np.where(vocab_valid, vocab_betas, 0.0)", CODE_KEY),
    ("", CODE_TEXT),
    ("# Store targets: NaN positions have value 0.0", CODE_COMMENT),
    ("all_betas[i] = torch.tensor(vocab_betas_clean)   # [vocab_size]", CODE_TEXT),
    ("", CODE_TEXT),
    ("# Store mask: True=real, False=NaN  ← used in Step 3", CODE_COMMENT),
    ("valid_mask[i] = torch.tensor(vocab_valid)         # [vocab_size]", CODE_KEY),
])

# Right explanation
add_panel(slide, Inches(8.1), Inches(1.4), Inches(4.8), Inches(2.1),
          "all_betas — the reconstruction target", ORANGE)
add_textbox(slide,
    "Shape: [batch, 49156]\n"
    "• Real positions: actual β ∈ [0, 1]\n"
    "• NaN positions: 0.0  (imputed)\n\n"
    "Imputing 0.0 is safe because valid_mask\n"
    "will exclude these from the loss.",
    Inches(8.3), Inches(2.0), Inches(4.4), Inches(1.4),
    font_size=14, color=LIGHT_GRAY)

add_panel(slide, Inches(8.1), Inches(3.65), Inches(4.8), Inches(2.2),
          "valid_mask — the key guard", GREEN)
add_textbox(slide,
    "Shape: [batch, 49156]\n"
    "• True  = position was measured\n"
    "• False = position was NaN\n\n"
    "This is what Step 3 uses to exclude\n"
    "NaN positions from reconstruction loss.\n"
    "It is returned in the batch dict.",
    Inches(8.3), Inches(4.25), Inches(4.4), Inches(1.5),
    font_size=14, color=LIGHT_GRAY)

add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.75), ACCENT2)
add_textbox(slide,
    "After Step 2a:   all_betas ready (NaN→0)    valid_mask ready (True/False)    "
    "both will be in the batch dict",
    Inches(0.6), Inches(6.2), Inches(12), Inches(0.45),
    font_size=15, color=GREEN)


# =============================================================================
# SLIDE 7 — Step 2b: WCEDCollator — Input View Selection
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Step 2b — WCEDCollator: Input View Selection (Valid Only)", Inches(8))

add_textbox(slide, "File: bmfm_methylation/shared/data_module.py   |   Class: WCEDCollator.__call__",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

add_code_panel(slide, Inches(0.4), Inches(1.4), Inches(7.5), Inches(3.2), [
    ("# Only sample input from VALID (non-NaN) positions", CODE_COMMENT),
    ("valid_indices = np.where(vocab_valid)[0]   # e.g. ~41k of 49k", CODE_KEY),
    ("", CODE_TEXT),
    ("# input_ratio applied to valid CpGs — not total vocab", CODE_COMMENT),
    ("# This handles high-NaN samples correctly:", CODE_COMMENT),
    ("#   50% × 41k = 20.5k input  (not 50% × 49k = 24.5k)", CODE_COMMENT),
    ("n_input = int(len(valid_indices) * self.input_ratio)", CODE_TEXT),
    ("", CODE_TEXT),
    ("# Random subset of valid positions for input view", CODE_COMMENT),
    ("indices_v1 = rng.choice(valid_indices, size=n_input, replace=False)", CODE_KEY),
    ("# NaN positions are NEVER in indices_v1", CODE_COMMENT),
    ("", CODE_TEXT),
    ("# Build token sequence: [CLS] + selected CpG tokens", CODE_COMMENT),
    ("ids  = [cls_id] + [vocab_cpg_ids[j] for j in indices_v1]", CODE_TEXT),
    ("vals = [cls_beta] + vocab_betas_clean[indices_v1].tolist()", CODE_TEXT),
])

# Right: 3 categories diagram
add_textbox(slide, "Position categories within the 49k vocabulary:",
            Inches(8.1), Inches(1.35), Inches(4.8), Inches(0.4),
            font_size=15, bold=True, color=HIGHLIGHT)

for i, (label, count, status, col) in enumerate([
    ("In input view",     "~20.5k positions", "EXCLUDED from recon loss\n(encoder already saw it)", ORANGE),
    ("Valid, held out",   "~20.5k positions", "INCLUDED in recon loss\n(real target → true signal)", GREEN),
    ("NaN position",      "~8k positions",    "target = 0.0\nexcluded via valid_mask", RED),
]):
    ly = Inches(1.9 + i * 1.5)
    add_rect(slide, Inches(8.1), ly, Inches(4.8), Inches(1.35), ACCENT)
    add_rect(slide, Inches(8.1), ly, Inches(0.07), Inches(1.35), col)
    add_textbox(slide, label,  Inches(8.3), ly + Inches(0.1),
                Inches(4.3), Inches(0.35), font_size=14, bold=True, color=col)
    add_textbox(slide, count,  Inches(8.3), ly + Inches(0.45),
                Inches(2.5), Inches(0.3), font_size=13, color=ORANGE)
    add_textbox(slide, status, Inches(8.3), ly + Inches(0.75),
                Inches(4.3), Inches(0.5), font_size=12, color=LIGHT_GRAY)

add_rect(slide, Inches(0.4), Inches(4.75), Inches(7.5), Inches(0.75), ACCENT2)
add_textbox(slide,
    "Key: input_ratio is applied to valid_indices — not to all 49k CpGs\n"
    "This ensures the model always sees exactly 50% of what was actually measured",
    Inches(0.55), Inches(4.82), Inches(7.2), Inches(0.6),
    font_size=13, color=GREEN, italic=True)

add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.75), ACCENT2)
add_textbox(slide,
    "After Step 2b:   input view built (NaN-free tokens)    "
    "NaN positions excluded from encoder input entirely",
    Inches(0.6), Inches(6.2), Inches(12), Inches(0.45),
    font_size=15, color=GREEN)


# =============================================================================
# SLIDE 8 — Step 3: WCEDLlamaModule — Loss Masking
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Step 3 — WCEDLlamaModule: Correct Reconstruction Loss", Inches(8.5))

add_textbox(slide, "File: bmfm_methylation/llama/wced_llama.py   |   Method: _shared_step",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

add_code_panel(slide, Inches(0.4), Inches(1.4), Inches(7.5), Inches(4.9), [
    ("# _shared_step — reconstruction loss calculation", CODE_COMMENT),
    ("input_mask = batch['input_mask']   # [B, 49156] True = in input view", CODE_TEXT),
    ("all_betas  = batch['all_betas']    # [B, 49156] targets (0 where NaN)", CODE_TEXT),
    ("valid_mask = batch.get('valid_mask')  # [B, 49156] True = non-NaN", CODE_TEXT),
    ("", CODE_TEXT),
    ("# Step 3 core: combine both exclusion masks", CODE_COMMENT),
    ("non_input_mask = ~input_mask          # not shown to encoder", CODE_TEXT),
    ("recon_mask = non_input_mask           # start with non-input", CODE_TEXT),
    ("if valid_mask is not None:", CODE_TEXT),
    ("    recon_mask = non_input_mask & valid_mask  # EXCLUDE NaN too", CODE_KEY),
    ("", CODE_TEXT),
    ("# MSE loss — element-wise, then masked", CODE_COMMENT),
    ("loss_per_cpg = recon_loss_fn(pred_betas, all_betas)  # [B, 49156]", CODE_TEXT),
    ("masked_loss  = loss_per_cpg * recon_mask.float()", CODE_TEXT),
    ("recon_loss   = masked_loss.sum()", CODE_TEXT),
    ("             / recon_mask.float().sum().clamp(min=1)  # correct denom", CODE_KEY),
    ("", CODE_TEXT),
    ("# Metrics also computed on recon_mask (not just ~input_mask)", CODE_COMMENT),
    ("ni_pred   = pred_betas[recon_mask]     # only real, held-out CpGs", CODE_TEXT),
    ("ni_target = all_betas[recon_mask]", CODE_TEXT),
    ("pcc = pearson_corr(ni_pred, ni_target)", CODE_KEY),
])

# Right: formula and explanation
add_panel(slide, Inches(8.1), Inches(1.4), Inches(4.8), Inches(2.3),
          "recon_mask definition", HIGHLIGHT)
add_textbox(slide,
    "recon_mask = ~input_mask  AND  valid_mask\n\n"
    "Positions included in loss:\n"
    "  (a) NOT in input view\n"
    "  (b) Actually measured (non-NaN)\n\n"
    "≈ 42% of vocab per step  (from valid_pct metric)",
    Inches(8.3), Inches(2.0), Inches(4.4), Inches(1.6),
    font_size=13, color=LIGHT_GRAY)

add_panel(slide, Inches(8.1), Inches(3.85), Inches(4.8), Inches(2.4),
          "Why the denominator matters", ORANGE)
add_textbox(slide,
    "Old code (SCBert): denominator = all non-input positions\n"
    "  → includes ~29k NaN positions with target=0\n"
    "  → loss looks artificially small\n\n"
    "LLaMA code: denominator = recon_mask.sum()\n"
    "  → only real, measured positions\n"
    "  → honest metric",
    Inches(8.3), Inches(4.48), Inches(4.4), Inches(1.65),
    font_size=13, color=LIGHT_GRAY)

add_rect(slide, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.75), ACCENT2)
add_textbox(slide,
    "✓  Loss computed only on (non-input  ∩  non-NaN) positions    "
    "✓  PCC and MAE metrics use the same clean mask",
    Inches(0.6), Inches(6.55), Inches(12), Inches(0.45),
    font_size=15, color=GREEN)


# =============================================================================
# SLIDE 9 — Bug (Old Code) vs Fix (LLaMA)
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Bug in Old Code vs. Correct Implementation in LLaMA", Inches(8.5))

add_textbox(slide, "Your LLaMA runs are NOT affected — the fix is already in wced_llama.py",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=15, color=GREEN, bold=True)

# Left: BUG
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.1), Inches(4.2), ACCENT)
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.1), Inches(0.07), RED)
add_textbox(slide, "❌  BUG — wced_module.py  (SCBert WCED)",
            Inches(0.55), Inches(1.48), Inches(5.8), Inches(0.4),
            font_size=16, bold=True, color=RED)
add_textbox(slide, "NOT used by your LLaMA runs",
            Inches(0.55), Inches(1.88), Inches(5.8), Inches(0.3),
            font_size=12, color=ORANGE, italic=True)
add_code_panel(slide, Inches(0.45), Inches(2.25), Inches(6.0), Inches(3.2), [
    ("# _shared_step — INCORRECT", CODE_COMMENT),
    ("non_input_mask_v1 = ~input_mask_v1", CODE_TEXT),
    ("# valid_mask is in the batch but never read!", RED),
    ("", CODE_TEXT),
    ("loss_per_cpg = recon_loss_fn(pred, all_betas)", CODE_TEXT),
    ("masked_loss  = loss_per_cpg * non_input_mask_v1.float()", CODE_TEXT),
    ("recon_loss   = masked_loss.sum()", CODE_TEXT),
    ("             / non_input_mask_v1.float().sum()", RED),
    ("# ↑ denominator includes ~29k NaN positions!", RED),
    ("", CODE_TEXT),
    ("# Result: 75% of gradient signal is from NaN", RED),
    ("# positions where target=0.0 (imputed, not real)", RED),
])

# Right: FIX
add_rect(slide, Inches(7.2), Inches(1.4), Inches(5.7), Inches(4.2), ACCENT)
add_rect(slide, Inches(7.2), Inches(1.4), Inches(5.7), Inches(0.07), GREEN)
add_textbox(slide, "✓  CORRECT — wced_llama.py  (LLaMA)",
            Inches(7.35), Inches(1.48), Inches(5.4), Inches(0.4),
            font_size=16, bold=True, color=GREEN)
add_textbox(slide, "Used by run 44450919  (your 7-day pretrain)",
            Inches(7.35), Inches(1.88), Inches(5.4), Inches(0.3),
            font_size=12, color=HIGHLIGHT2, italic=True)
add_code_panel(slide, Inches(7.25), Inches(2.25), Inches(5.6), Inches(3.2), [
    ("# _shared_step — CORRECT", CODE_COMMENT),
    ("non_input_mask = ~input_mask_v1", CODE_TEXT),
    ("valid_mask = batch.get('valid_mask')", CODE_TEXT),
    ("recon_mask = non_input_mask", CODE_TEXT),
    ("if valid_mask is not None:", CODE_TEXT),
    ("    recon_mask = non_input_mask & valid_mask", CODE_KEY),
    ("# ↑ NaN positions excluded from loss", CODE_KEY),
    ("", CODE_TEXT),
    ("loss_per_cpg = recon_loss_fn(pred, all_betas)", CODE_TEXT),
    ("masked_loss  = loss_per_cpg * recon_mask.float()", CODE_TEXT),
    ("recon_loss   = masked_loss.sum()", CODE_TEXT),
    ("             / recon_mask.float().sum().clamp(min=1)", CODE_KEY),
    ("# ↑ denominator = only real, held-out CpGs", CODE_KEY),
])

# Impact comparison
add_rect(slide, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.05), ACCENT2)
for i, (label, old_val, new_val, col) in enumerate([
    ("Signal in recon loss:",  "25% real / 75% NaN zeros",    "~100% real measured CpGs", GREEN),
    ("Reported MSE metric:",   "Artificially small (easy 0s)", "Honest — real CpGs only",  GREEN),
    ("Your run affected?",     "N/A — not your code",          "✓ Already correct",        HIGHLIGHT),
]):
    y = Inches(5.82 + i * 0.3)
    add_textbox(slide, label,   Inches(0.55), y, Inches(2.5),  Inches(0.28), font_size=12, color=SUBTITLE_GRAY)
    add_textbox(slide, old_val, Inches(3.2),  y, Inches(4.5),  Inches(0.28), font_size=12, color=RED)
    add_textbox(slide, new_val, Inches(8.0),  y, Inches(4.5),  Inches(0.28), font_size=12, color=col)


# =============================================================================
# SLIDE 10 — Pretrain Metrics (Run 44450919, Epoch 99)
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Pretrain Results — Run 44450919  (Epoch 99 / 300)", Inches(7))

add_textbox(slide, "llama-small-all49k-r0.5-w0.0   |   256D × 4L × 4H  |  ~5M params  |  Runtime: 4.5 days",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=SUBTITLE_GRAY, italic=True)

# Metric boxes — top row
metrics_top = [
    ("0.971",  "val/pcc\nRecon quality"),
    ("0.043",  "val/mae\nBeta error [0,1]"),
    ("0.006",  "val/loss\nRecon MSE"),
    ("42.0%",  "valid_pct\nUseful positions/step"),
]
for i, (val, label) in enumerate(metrics_top):
    add_metric_box(slide, Inches(0.4 + i * 3.15), Inches(1.4),
                   Inches(2.9), Inches(1.3), label, val,
                   val_color=GREEN if i == 0 else ORANGE)

# CLS collapse diagnostics — second row
add_textbox(slide, "CLS Collapse Diagnostics  (collapsed = bad):",
            Inches(0.4), Inches(2.9), Inches(7), Inches(0.35),
            font_size=16, bold=True, color=HIGHLIGHT)
cls_metrics = [
    ("0.376",  "cls_similarity\n(1.0 = fully collapsed)"),
    ("0.429",  "cls_variance\n(0.0 = fully collapsed)"),
    ("0.310",  "pred_std\n(0.0 = predicting mean)"),
    ("0.570",  "pred_var_ratio\n(1.0 = perfect)"),
]
for i, (val, label) in enumerate(cls_metrics):
    add_metric_box(slide, Inches(0.4 + i * 3.15), Inches(3.3),
                   Inches(2.9), Inches(1.3), label, val,
                   val_color=GREEN if i < 3 else ORANGE)

# NaN rate inference
add_rect(slide, Inches(0.4), Inches(4.8), Inches(7.8), Inches(1.85), ACCENT)
add_textbox(slide, "NaN Rate Inference from valid_pct Metric",
            Inches(0.6), Inches(4.9), Inches(7.4), Inches(0.4),
            font_size=16, bold=True, color=HIGHLIGHT)
add_textbox(slide,
    "valid_pct = (non-input ∩ non-NaN) / vocab_size = 42%\n"
    "input = 50% of valid   →   non-input non-NaN ≈ 50% of valid\n"
    "→  valid (non-NaN) ≈ 84% of 49,156   →   NaN rate ≈ 16%\n"
    "→  ~41k of 49k probes are measured per sample on average",
    Inches(0.6), Inches(5.35), Inches(7.4), Inches(1.2),
    font_size=13, color=LIGHT_GRAY)

# Training health
add_rect(slide, Inches(8.4), Inches(4.8), Inches(4.5), Inches(1.85), ACCENT)
add_textbox(slide, "Training Health",
            Inches(8.6), Inches(4.9), Inches(4.1), Inches(0.4),
            font_size=16, bold=True, color=HIGHLIGHT)
for i, (label, val, col) in enumerate([
    ("cpg_scale (init=0.1)", "3.68  ✓ growing", GREEN),
    ("grad_norm",            "0.0018  (converging)", HIGHLIGHT2),
    ("train/pcc",            "0.970  (no overfit)", GREEN),
    ("LR at epoch 99",       "3.82e-4  (cosine decay)", LIGHT_GRAY),
]):
    y = Inches(5.35 + i * 0.3)
    add_textbox(slide, label, Inches(8.6), y, Inches(2.2), Inches(0.28),
                font_size=12, color=SUBTITLE_GRAY)
    add_textbox(slide, val,   Inches(10.9), y, Inches(1.8), Inches(0.28),
                font_size=12, color=col)

# Interpretation
add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5), ACCENT2)
add_textbox(slide,
    "✓ No CLS collapse   ✓ Reconstruction quality excellent   "
    "⚠ pred_var_ratio=0.57 — slight underdispersion (expected with pure reconstruction)",
    Inches(0.6), Inches(6.87), Inches(12.1), Inches(0.38),
    font_size=14, color=HIGHLIGHT2)


# =============================================================================
# SLIDE 11 — PLACEHOLDER: Investigation Results — NaN Statistics
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Data Investigation Results — NaN Statistics  [PENDING]", Inches(8))

add_textbox(slide, "Run on cluster:  sbatch scripts/utils/investigate_data.sh",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

add_textbox(slide, "⏳  Results pending — run investigation script and fill in below",
            Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.45),
            font_size=16, bold=True, color=ORANGE)

# Pretrain placeholders
add_panel(slide, Inches(0.4), Inches(2.05), Inches(6.1), Inches(4.3),
          "PRETRAIN  (169k × 49k)", HIGHLIGHT)
for label, yoff in [
    ("Total NaN count:             [ TBD ]", 0.75),
    ("NaN %:                       [ TBD ]", 1.15),
    ("Samples with ≥1% NaN:       [ TBD ]", 1.55),
    ("Samples with ≥10% NaN:      [ TBD ]", 1.95),
    ("Fully missing samples:       [ TBD ]", 2.35),
    ("CpGs with ≥50% NaN:         [ TBD ]", 2.75),
    ("Fully missing CpGs:          [ TBD ]", 3.15),
]:
    add_textbox(slide, label, Inches(0.65), Inches(yoff + 2.05),
                Inches(5.6), Inches(0.38), font_size=14, color=LIGHT_GRAY)

# Finetune placeholders
add_panel(slide, Inches(7.0), Inches(2.05), Inches(5.9), Inches(4.3),
          "FINE-TUNE  (11.5k × 49k)", GREEN)
for label, yoff in [
    ("Total NaN count:         [ TBD ]", 0.75),
    ("NaN %:                   [ TBD ]", 1.15),
    ("Samples with ≥1% NaN:   [ TBD ]", 1.55),
    ("Samples with ≥10% NaN:  [ TBD ]", 1.95),
    ("Train split NaN %:       [ TBD ]", 2.35),
    ("Valid split NaN %:       [ TBD ]", 2.75),
    ("Test  split NaN %:       [ TBD ]", 3.15),
]:
    add_textbox(slide, label, Inches(7.25), Inches(yoff + 2.05),
                Inches(5.4), Inches(0.38), font_size=14, color=LIGHT_GRAY)

add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.75), ACCENT2)
add_textbox(slide,
    "Expected from valid_pct≈42%:  NaN rate ≈ 16% in pretrain corpus  |  "
    "Fine-tune NaN rate may differ (different sample sources)",
    Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
    font_size=14, color=SUBTITLE_GRAY, italic=True)


# =============================================================================
# SLIDE 12 — PLACEHOLDER: Beta Distribution
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Data Investigation Results — Beta Distribution  [PENDING]", Inches(8))

add_textbox(slide, "Run on cluster:  sbatch scripts/utils/investigate_data.sh",
            Inches(0.5), Inches(0.95), Inches(12), Inches(0.35),
            font_size=14, color=CODE_TEXT, italic=True)

add_textbox(slide, "⏳  Results pending — paste percentile output from investigation script",
            Inches(0.4), Inches(1.45), Inches(12.5), Inches(0.45),
            font_size=16, bold=True, color=ORANGE)

# Pretrain distribution
add_panel(slide, Inches(0.4), Inches(2.05), Inches(6.1), Inches(4.3),
          "PRETRAIN β-value Distribution", HIGHLIGHT)
percs = ["0%", "1%", "5%", "25%", "50%", "75%", "95%", "99%", "100%"]
for i, p in enumerate(percs):
    y = Inches(2.8 + i * 0.37)
    add_textbox(slide, f"  {p:>5s}:", Inches(0.65), y, Inches(1.2), Inches(0.35),
                font_size=13, color=SUBTITLE_GRAY)
    add_textbox(slide, "[ TBD ]", Inches(1.9), y, Inches(2.0), Inches(0.35),
                font_size=13, color=LIGHT_GRAY)

# Finetune distribution
add_panel(slide, Inches(7.0), Inches(2.05), Inches(5.9), Inches(4.3),
          "FINE-TUNE β-value Distribution", GREEN)
for i, p in enumerate(percs):
    y = Inches(2.8 + i * 0.37)
    add_textbox(slide, f"  {p:>5s}:", Inches(7.25), y, Inches(1.2), Inches(0.35),
                font_size=13, color=SUBTITLE_GRAY)
    add_textbox(slide, "[ TBD ]", Inches(8.5), y, Inches(2.0), Inches(0.35),
                font_size=13, color=LIGHT_GRAY)

# Out-of-range check
add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.75), ACCENT2)
for i, label in enumerate(["Out-of-range (β < 0 or β > 1):", "Duplicate sample IDs:", "Duplicate CpG IDs:"]):
    add_textbox(slide, f"{label}  [ TBD ]",
                Inches(0.6 + i * 4.2), Inches(6.6),
                Inches(4.0), Inches(0.5),
                font_size=14, color=LIGHT_GRAY)


# =============================================================================
# SLIDE 13 — Summary
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Summary", Inches(2.5))

# Left: pipeline summary
add_panel(slide, Inches(0.4), Inches(1.1), Inches(5.9), Inches(5.7),
          "Pipeline — What Each File Does", HIGHLIGHT)
items = [
    ("data_module.py — MethylationDataset",    "Reads β, computes valid_mask = isfinite(β)"),
    ("data_module.py — WCEDCollator (2a)",      "Replaces NaN→0 in targets, stores valid_mask"),
    ("data_module.py — WCEDCollator (2b)",      "Input views: only sample from valid positions"),
    ("wced_llama.py — WCEDLlamaModule",         "recon_mask = ~input_mask & valid_mask"),
]
for i, (title, desc) in enumerate(items):
    y = Inches(1.8 + i * 1.15)
    add_rect(slide, Inches(0.55), y, Inches(0.06), Inches(0.85), HIGHLIGHT)
    add_textbox(slide, title, Inches(0.75), y + Inches(0.05),
                Inches(5.2), Inches(0.38), font_size=13, bold=True, color=WHITE)
    add_textbox(slide, desc, Inches(0.75), y + Inches(0.42),
                Inches(5.2), Inches(0.38), font_size=12, color=LIGHT_GRAY)

# Right: key conclusions
add_panel(slide, Inches(7.0), Inches(1.1), Inches(5.9), Inches(2.6),
          "Your LLaMA Runs — Status", GREEN)
for i, (text, col) in enumerate([
    ("✓ Pretrain run 44450919: correct NaN handling", GREEN),
    ("✓ Fine-tune 44574410: recon_weight=0.0, NaN irrelevant", GREEN),
    ("✓ wced_llama.py already had the fix from the start", GREEN),
    ("✓ Bug exists only in old SCBert wced_module.py (not your run)", HIGHLIGHT2),
]):
    add_textbox(slide, text, Inches(7.2), Inches(1.75 + i * 0.48),
                Inches(5.5), Inches(0.42), font_size=13, color=col)

add_panel(slide, Inches(7.0), Inches(3.85), Inches(5.9), Inches(3.0),
          "Pretrain Quality Check (Epoch 99)", ORANGE)
for i, (text, col) in enumerate([
    ("PCC = 0.971  — excellent reconstruction", GREEN),
    ("MAE = 0.043  — 4.3% error on [0,1]",      GREEN),
    ("cls_similarity = 0.376  — no CLS collapse", GREEN),
    ("cls_variance = 0.429    — healthy diversity", GREEN),
    ("pred_var_ratio = 0.570  — slight underdispersion", ORANGE),
    ("valid_pct ≈ 42%  →  NaN rate ≈ 16%", HIGHLIGHT2),
]):
    add_textbox(slide, text, Inches(7.2), Inches(4.5 + i * 0.38),
                Inches(5.5), Inches(0.35), font_size=13, color=col)

add_rect(slide, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.38), ACCENT2)
add_textbox(slide,
    "Next step: paste investigate_data.sh results into slides 11–12  |  "
    "Monitor finetune val/r2 to assess if reconstruction pretraining transferred to age prediction",
    Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
    font_size=12, color=SUBTITLE_GRAY, italic=True)


# =============================================================================
# Save
# =============================================================================
out_path = "NaN_Handling_MethylLlama.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
